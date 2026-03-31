"""Office document converters.

Supported conversions (require optional dependencies):
  docx → html, md, txt, pdf
  pdf  → txt, md
  xlsx → csv, json, yaml
  csv  → xlsx
  pptx → txt, html
"""

from __future__ import annotations

import sys
from pathlib import Path

from cvt.backends.base import BaseConverter, ConversionError, _try_import
from cvt.log import get_logger

log = get_logger(__name__)


# ---------------------------------------------------------------------------
# DOCX → HTML / Markdown / TXT  (mammoth)
# ---------------------------------------------------------------------------

class DocxToHtmlConverter(BaseConverter):
    name = "mammoth"
    supported = [("docx", "html")]
    requires = ["mammoth"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import mammoth  # type: ignore[import]
        with open(src, "rb") as fh:
            result = mammoth.convert_to_html(fh)
        if result.messages:
            for msg in result.messages:
                log.debug("mammoth: %s", msg)
        html = f"<!DOCTYPE html>\n<html><head><meta charset='utf-8'></head><body>\n{result.value}\n</body></html>"
        dst.write_text(html, encoding="utf-8")
        log.debug("Converted %s → %s (docx→html)", src.name, dst.name)


class DocxToMarkdownConverter(BaseConverter):
    name = "mammoth-md"
    supported = [("docx", "md")]
    requires = ["mammoth"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import mammoth  # type: ignore[import]
        with open(src, "rb") as fh:
            result = mammoth.convert_to_markdown(fh)
        dst.write_text(result.value, encoding="utf-8")
        log.debug("Converted %s → %s (docx→md)", src.name, dst.name)


class DocxToTxtConverter(BaseConverter):
    name = "mammoth-txt"
    supported = [("docx", "txt")]
    requires = ["mammoth"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import mammoth  # type: ignore[import]
        with open(src, "rb") as fh:
            result = mammoth.extract_raw_text(fh)
        dst.write_text(result.value, encoding="utf-8")
        log.debug("Converted %s → %s (docx→txt)", src.name, dst.name)


# ---------------------------------------------------------------------------
# DOCX → PDF  (docx2pdf on Win/Mac, or via python-docx + weasyprint on Linux)
# ---------------------------------------------------------------------------

class DocxToPdfConverter(BaseConverter):
    name = "docx2pdf"
    supported = [("docx", "pdf")]

    @classmethod
    def is_available(cls) -> bool:
        if sys.platform in ("win32", "darwin"):
            return _try_import("docx2pdf")
        # On Linux: use mammoth + weasyprint as fallback
        return _try_import("mammoth") and _try_import("weasyprint")

    requires = []  # checked dynamically above

    def convert(self, src: Path, dst: Path, **options) -> None:
        if sys.platform in ("win32", "darwin"):
            try:
                import docx2pdf  # type: ignore[import]
                docx2pdf.convert(str(src), str(dst))
            except Exception as exc:
                raise ConversionError(f"docx2pdf failed: {exc}") from exc
        else:
            # Linux fallback: docx → html → pdf via weasyprint
            import mammoth  # type: ignore[import]
            from weasyprint import HTML  # type: ignore[import]
            with open(src, "rb") as fh:
                result = mammoth.convert_to_html(fh)
            html = f"<!DOCTYPE html>\n<html><head><meta charset='utf-8'><style>body{{font-family:sans-serif;max-width:800px;margin:auto;padding:2em}}</style></head><body>\n{result.value}\n</body></html>"
            try:
                HTML(string=html).write_pdf(str(dst))
            except Exception as exc:
                raise ConversionError(f"PDF generation failed: {exc}") from exc
        log.debug("Converted %s → %s (docx→pdf)", src.name, dst.name)


# ---------------------------------------------------------------------------
# PDF → TXT  (pdfminer.six)
# ---------------------------------------------------------------------------

class PdfToTxtConverter(BaseConverter):
    name = "pdfminer"
    supported = [("pdf", "txt")]
    requires = ["pdfminer"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        from pdfminer.high_level import extract_text  # type: ignore[import]
        try:
            text = extract_text(str(src))
        except Exception as exc:
            raise ConversionError(f"PDF text extraction failed: {exc}") from exc
        dst.write_text(text or "", encoding="utf-8")
        log.debug("Converted %s → %s (pdf→txt)", src.name, dst.name)


# ---------------------------------------------------------------------------
# PDF → Markdown  (pdfminer → lightweight md wrapping)
# ---------------------------------------------------------------------------

class PdfToMarkdownConverter(BaseConverter):
    name = "pdfminer-md"
    supported = [("pdf", "md")]
    requires = ["pdfminer"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        from pdfminer.high_level import extract_text  # type: ignore[import]
        try:
            text = extract_text(str(src)) or ""
        except Exception as exc:
            raise ConversionError(f"PDF text extraction failed: {exc}") from exc
        # Wrap pages/paragraphs in minimal markdown
        lines = text.splitlines()
        md_lines = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                md_lines.append("")
            else:
                md_lines.append(stripped)
        dst.write_text("\n".join(md_lines), encoding="utf-8")
        log.debug("Converted %s → %s (pdf→md)", src.name, dst.name)


# ---------------------------------------------------------------------------
# XLSX → CSV / JSON / YAML  (openpyxl)
# ---------------------------------------------------------------------------

class XlsxToCsvConverter(BaseConverter):
    name = "openpyxl-csv"
    supported = [("xlsx", "csv"), ("xls", "csv")]
    requires = ["openpyxl"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import csv
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(str(src), read_only=True, data_only=True)
        ws = wb.active
        with open(dst, "w", newline="", encoding="utf-8") as fh:
            writer = csv.writer(fh)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([str(c) if c is not None else "" for c in row])
        wb.close()
        log.debug("Converted %s → %s (xlsx→csv)", src.name, dst.name)


class XlsxToJsonConverter(BaseConverter):
    name = "openpyxl-json"
    supported = [("xlsx", "json")]
    requires = ["openpyxl"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import json
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(str(src), read_only=True, data_only=True)
        ws = wb.active
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            dst.write_text("[]", encoding="utf-8")
            return
        headers = [str(c) if c is not None else f"col_{i}" for i, c in enumerate(rows[0])]
        data = []
        for row in rows[1:]:
            data.append({headers[i]: row[i] for i in range(len(headers))})
        with open(dst, "w", encoding="utf-8") as fh:
            json.dump(data, fh, ensure_ascii=False, indent=2, default=str)
        wb.close()
        log.debug("Converted %s → %s (xlsx→json)", src.name, dst.name)


# ---------------------------------------------------------------------------
# CSV → XLSX  (openpyxl)
# ---------------------------------------------------------------------------

class CsvToXlsxConverter(BaseConverter):
    name = "openpyxl-xlsx"
    supported = [("csv", "xlsx")]
    requires = ["openpyxl"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import csv
        import openpyxl  # type: ignore[import]
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(src, newline="", encoding="utf-8") as fh:
            for row in csv.reader(fh):
                ws.append(row)
        wb.save(str(dst))
        log.debug("Converted %s → %s (csv→xlsx)", src.name, dst.name)


# ---------------------------------------------------------------------------
# PPTX → TXT / HTML  (python-pptx)
# ---------------------------------------------------------------------------

class PptxToTxtConverter(BaseConverter):
    name = "pptx-txt"
    supported = [("pptx", "txt")]
    requires = ["pptx"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(str(src))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            lines.append("")
        dst.write_text("\n".join(lines), encoding="utf-8")
        log.debug("Converted %s → %s (pptx→txt)", src.name, dst.name)


class PptxToHtmlConverter(BaseConverter):
    name = "pptx-html"
    supported = [("pptx", "html")]
    requires = ["pptx"]

    def convert(self, src: Path, dst: Path, **options) -> None:
        import html as html_mod
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(str(src))
        sections = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"<h2>Slide {i}</h2>"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(f"<p>{html_mod.escape(text)}</p>")
            sections.append("\n".join(parts))
        body = "\n<hr/>\n".join(sections)
        full = f"<!DOCTYPE html>\n<html><head><meta charset='utf-8'></head><body>\n{body}\n</body></html>"
        dst.write_text(full, encoding="utf-8")
        log.debug("Converted %s → %s (pptx→html)", src.name, dst.name)
